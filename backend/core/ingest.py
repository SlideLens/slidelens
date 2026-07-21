"""Deck / audio ingest (render slides to PNG, extract payloads).

Only this module shells out to ``soffice``/``ffmpeg``; ``_run_subprocess`` is the
single seam, so unit tests mock one function instead of the OS process layer.
"""

from __future__ import annotations

import asyncio
from collections.abc import Awaitable, Callable
from pathlib import Path

from pdf2image import convert_from_path
from pdf2image.exceptions import PDFPageCountError, PDFPopplerTimeoutError, PDFSyntaxError
from pptx import Presentation

from core.constants import (
    DECK_RENDER_TIMEOUT_SECONDS,
    MAX_AUDIO_MINUTES,
    MAX_DECK_SIZE_MB,
    MAX_DECK_SLIDES,
    SLIDE_PNG_DPI,
    WAV_16K_MONO_BYTES_PER_SECOND,
    AllowedDeckFormat,
)
from core.context import ReviewContext
from core.ingest_errors import (
    AudioTooLongError,
    CorruptedFileError,
    DeckTooLargeError,
    EmptyDeckError,
    NoAudioTrackError,
    PasswordProtectedError,
    RenderTimeoutError,
    UnsupportedDeckFormatError,
)

RunSubprocess = Callable[..., Awaitable[tuple[int, bytes, bytes]]]


async def _run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
    """Run ``cmd``, killing it on timeout. The default seam for ingest subprocess calls."""
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    try:
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
    except TimeoutError:
        proc.kill()
        await proc.wait()
        raise
    return proc.returncode or 0, stdout, stderr


def _resolve_deck_format(path: Path) -> AllowedDeckFormat:
    suffix = path.suffix.lower()
    try:
        return AllowedDeckFormat(suffix)
    except ValueError:
        allowed = ", ".join(sorted(fmt.value for fmt in AllowedDeckFormat))
        raise UnsupportedDeckFormatError(
            f"Unsupported deck format {suffix!r}; allowed: {allowed}: {path}"
        ) from None


class DeckIngestor:
    """PPTX/PDF → numbered slide PNGs + per-slide text."""

    def __init__(
        self,
        *,
        soffice_bin: str = "soffice",
        run_subprocess: RunSubprocess = _run_subprocess,
        timeout_seconds: float = DECK_RENDER_TIMEOUT_SECONDS,
        dpi: int = SLIDE_PNG_DPI,
    ) -> None:
        self._soffice_bin = soffice_bin
        self._run_subprocess = run_subprocess
        self._timeout = timeout_seconds
        self._dpi = dpi

    async def ingest(
        self,
        deck: Path,
        workdir: Path,
        ctx: ReviewContext | None = None,
    ) -> list[Path]:
        if not deck.is_file():
            raise CorruptedFileError(f"Deck file not found: {deck}")

        size_mb = deck.stat().st_size / (1024 * 1024)
        if size_mb > MAX_DECK_SIZE_MB:
            raise DeckTooLargeError(
                f"Deck is {size_mb:.1f} MB, limit is {MAX_DECK_SIZE_MB} MB: {deck}"
            )

        workdir.mkdir(parents=True, exist_ok=True)
        deck_format = _resolve_deck_format(deck)
        slide_texts: dict[int, str] = {}

        if deck_format is AllowedDeckFormat.PPTX:
            slide_texts = self._extract_pptx_texts(deck)
            if not slide_texts:
                raise EmptyDeckError(f"Deck has no slides: {deck}")
            if len(slide_texts) > MAX_DECK_SLIDES:
                raise DeckTooLargeError(
                    f"Deck has {len(slide_texts)} slides, limit is {MAX_DECK_SLIDES}: {deck}"
                )
            pdf_path = await self._convert_to_pdf(deck, workdir)
        else:
            pdf_path = deck

        pngs = self._render_pdf_to_pngs(pdf_path, workdir)
        if not pngs:
            raise EmptyDeckError(f"Rendered zero pages from: {deck}")
        if len(pngs) > MAX_DECK_SLIDES:
            raise DeckTooLargeError(
                f"Deck rendered {len(pngs)} pages, limit is {MAX_DECK_SLIDES}: {deck}"
            )

        if ctx is not None:
            if slide_texts:
                ctx.meta["slide_texts"] = slide_texts
            for idx, png in enumerate(pngs, start=1):
                ctx.slide_pngs[idx] = png

        return pngs

    def _extract_pptx_texts(self, deck: Path) -> dict[int, str]:
        try:
            presentation = Presentation(deck)
        except Exception as exc:
            raise CorruptedFileError(f"Deck is not a valid PPTX: {deck}") from exc

        texts: dict[int, str] = {}
        for idx, slide in enumerate(presentation.slides, start=1):
            chunks = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            texts[idx] = "\n".join(chunks)
        return texts

    async def _convert_to_pdf(self, deck: Path, workdir: Path) -> Path:
        cmd = [
            self._soffice_bin,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(workdir),
            str(deck),
        ]
        last_error = "unknown error"
        for _attempt in range(2):
            try:
                returncode, _stdout, stderr = await self._run_subprocess(
                    cmd, timeout=self._timeout
                )
            except TimeoutError:
                last_error = f"soffice timed out after {self._timeout}s"
                continue
            if returncode == 0:
                pdf_path = workdir / f"{deck.stem}.pdf"
                if pdf_path.is_file():
                    return pdf_path
                last_error = f"soffice exited 0 but {pdf_path.name} is missing"
                continue
            last_error = f"soffice exited {returncode}: {stderr.decode(errors='replace')[:500]}"
        raise RenderTimeoutError(
            f"PPTX→PDF render of {deck} failed after 2 attempts: {last_error}"
        )

    def _render_pdf_to_pngs(self, pdf_path: Path, workdir: Path) -> list[Path]:
        try:
            images = convert_from_path(str(pdf_path), dpi=self._dpi)
        except PDFPopplerTimeoutError as exc:
            raise RenderTimeoutError(str(exc)) from exc
        except (PDFPageCountError, PDFSyntaxError) as exc:
            message = str(exc).lower()
            if "password" in message or "encrypt" in message:
                raise PasswordProtectedError(str(exc)) from exc
            raise CorruptedFileError(str(exc)) from exc

        paths: list[Path] = []
        for idx, image in enumerate(images, start=1):
            png_path = workdir / f"slide_{idx:03d}.png"
            image.save(png_path, "PNG")
            paths.append(png_path)
        return paths


class AudioExtractor:
    """Video/audio → 16kHz mono WAV via ``ffmpeg``."""

    def __init__(
        self,
        *,
        ffmpeg_bin: str = "ffmpeg",
        run_subprocess: RunSubprocess = _run_subprocess,
        timeout_seconds: float = 120.0,
        max_minutes: float = MAX_AUDIO_MINUTES,
    ) -> None:
        self._ffmpeg_bin = ffmpeg_bin
        self._run_subprocess = run_subprocess
        self._timeout = timeout_seconds
        self._max_minutes = max_minutes

    async def extract(self, media: Path, workdir: Path) -> Path:
        if not media.is_file():
            raise CorruptedFileError(f"Media file not found: {media}")

        workdir.mkdir(parents=True, exist_ok=True)
        out_path = workdir / f"{media.stem}.wav"
        cmd = [
            self._ffmpeg_bin,
            "-y",
            "-i",
            str(media),
            "-vn",
            "-ac",
            "1",
            "-ar",
            "16000",
            str(out_path),
        ]
        try:
            returncode, _stdout, stderr = await self._run_subprocess(cmd, timeout=self._timeout)
        except TimeoutError as exc:
            raise RenderTimeoutError(f"ffmpeg timed out on {media}") from exc

        text = stderr.decode(errors="replace").lower()
        if returncode != 0:
            if "does not contain any stream" in text or "stream map '0:a" in text:
                raise NoAudioTrackError(f"No audio track in {media}")
            raise CorruptedFileError(f"ffmpeg failed on {media}: {text[:500]}")
        if not out_path.is_file() or out_path.stat().st_size == 0:
            raise NoAudioTrackError(f"No audio track in {media}")

        # Потолок ставим по факту извлечения: длительность из размера WAV точна,
        # а размер исходника о ней почти ничего не говорит (битрейт гуляет в разы).
        duration_seconds = out_path.stat().st_size / WAV_16K_MONO_BYTES_PER_SECOND
        if duration_seconds > self._max_minutes * 60:
            out_path.unlink(missing_ok=True)
            raise AudioTooLongError(
                f"Запись питча длиннее {self._max_minutes:.0f} минут "
                f"({duration_seconds / 60:.0f} мин)"
            )
        return out_path
