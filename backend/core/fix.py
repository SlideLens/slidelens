"""PPTX auto-fix strategies (safe rule set only, ADR 0006).

``PptxFixer`` copies the deck, applies each ``auto_fixable`` finding's matching
``FixRule`` to the shape closest to its bbox, then runs a structural control
check and rolls back to the original if a touched shape ended up off-slide.
Slide relayout is deliberately out of scope.
"""

from __future__ import annotations

import json
import re
import shutil
from abc import ABC, abstractmethod
from pathlib import Path
from typing import ClassVar

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.util import Pt

from core.fix_schemas import FixLogEntry, FixResult
from core.geometry import iou
from core.schemas import BBox, Category, Finding

MIN_FONT_SIZE_PT = 14
CONTRAST_MIN_RATIO = 4.5
EMU_PER_PX = 9525
ALIGNMENT_TOLERANCE_EMU = 3 * EMU_PER_PX
FIXED_DECK_FILENAME = "fixed.pptx"
_UNSAFE_FILENAME_CHARS = re.compile(r'[\\/:*?"<>|]')


def fixed_deck_filename(original_filename: str, version: int) -> str:
    """Public download name for the Исправленная дека: ``{дека}_исправленный_версия№{N}.pptx``."""
    stem = _UNSAFE_FILENAME_CHARS.sub("_", Path(original_filename).stem).strip() or "Дека"
    return f"{stem}_исправленная_версия№{version}.pptx"


class FixRule(ABC):
    """One safe, narrow PPTX edit strategy."""

    name: ClassVar[str]

    @abstractmethod
    def applies_to(self, finding: Finding) -> bool:
        """Whether this rule is the right one for this finding."""

    @abstractmethod
    def apply(self, shape: object, finding: Finding, slide: Slide) -> FixResult:
        """Attempt the edit; ``slide`` gives access to sibling shapes."""


def _relative_luminance(rgb: tuple[int, int, int]) -> float:
    def channel(value: int) -> float:
        c = value / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = (channel(v) for v in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    la, lb = _relative_luminance(a), _relative_luminance(b)
    lighter, darker = max(la, lb), min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


class MinFontSizeRule(FixRule):
    """Font runs smaller than 14pt are raised to 14pt."""

    name = "min_font_size"

    def applies_to(self, finding: Finding) -> bool:
        if finding.category != Category.TYPOGRAPHY:
            return False
        text = f"{finding.title} {finding.description}".lower()
        return "шрифт" in text or "кегл" in text

    def apply(self, shape: object, finding: Finding, slide: Slide) -> FixResult:
        if not getattr(shape, "has_text_frame", False):
            return FixResult(applied=False, reason="shape has no text frame")
        changed = False
        for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
            for run in paragraph.runs:
                if run.font.size is not None and run.font.size.pt < MIN_FONT_SIZE_PT:
                    run.font.size = Pt(MIN_FONT_SIZE_PT)
                    changed = True
        return FixResult(applied=changed, reason="" if changed else "no run below minimum size")


class ContrastRule(FixRule):
    """Text below the WCAG 4.5:1 ratio is recolored to black or white."""

    name = "contrast"

    def applies_to(self, finding: Finding) -> bool:
        if finding.category != Category.READABILITY:
            return False
        return "контраст" in f"{finding.title} {finding.description}".lower()

    def apply(self, shape: object, finding: Finding, slide: Slide) -> FixResult:
        if not getattr(shape, "has_text_frame", False):
            return FixResult(applied=False, reason="shape has no text frame")
        background = self._background_color(shape)
        changed = False
        for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
            for run in paragraph.runs:
                color = self._run_color(run)
                if color is None:
                    continue
                if _contrast_ratio(color, background) < CONTRAST_MIN_RATIO:
                    is_bg_light = _relative_luminance(background) > 0.5
                    new_color = (0, 0, 0) if is_bg_light else (255, 255, 255)
                    run.font.color.rgb = RGBColor(*new_color)
                    changed = True
        reason = "" if changed else "no run below contrast threshold"
        return FixResult(applied=changed, reason=reason)

    @staticmethod
    def _background_color(shape: object) -> tuple[int, int, int]:
        try:
            return tuple(shape.fill.fore_color.rgb)  # type: ignore[attr-defined]
        except (AttributeError, TypeError):
            return (255, 255, 255)

    @staticmethod
    def _run_color(run: object) -> tuple[int, int, int] | None:
        try:
            return tuple(run.font.color.rgb)  # type: ignore[attr-defined]
        except (AttributeError, TypeError):
            return None


class AlignmentRule(FixRule):
    """A shape's edge within 3px of a sibling's snaps to match it exactly."""

    name = "alignment"

    def applies_to(self, finding: Finding) -> bool:
        text = f"{finding.title} {finding.description}".lower()
        return "выровн" in text or "выравн" in text

    def apply(self, shape: object, finding: Finding, slide: Slide) -> FixResult:
        changed = False
        for attr in ("left", "top"):
            target = getattr(shape, attr)
            if target is None:
                continue
            for other in slide.shapes:
                if other is shape:
                    continue
                other_value = getattr(other, attr)
                if other_value is None or other_value == target:
                    continue
                if abs(other_value - target) <= ALIGNMENT_TOLERANCE_EMU:
                    setattr(shape, attr, other_value)
                    changed = True
                    break
        return FixResult(applied=changed, reason="" if changed else "no near-aligned sibling found")


def _find_matching_shape(
    slide: Slide, bbox: BBox, slide_width: int, slide_height: int
) -> object | None:
    best_shape: object | None = None
    best_iou = 0.0
    for shape in slide.shapes:
        if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
            continue
        shape_bbox = BBox(
            x=max(0.0, min(1.0, shape.left / slide_width)),
            y=max(0.0, min(1.0, shape.top / slide_height)),
            w=max(0.0, min(1.0, shape.width / slide_width)),
            h=max(0.0, min(1.0, shape.height / slide_height)),
        )
        score = iou(bbox, shape_bbox)
        if score > best_iou:
            best_iou = score
            best_shape = shape
    return best_shape if best_iou > 0 else None


class PptxFixer:
    """Copies the deck, applies safe rules to auto_fixable findings, logs + rolls back."""

    def __init__(self, rules: list[FixRule] | None = None) -> None:
        self._rules = rules or [MinFontSizeRule(), ContrastRule(), AlignmentRule()]
        self.last_log: list[FixLogEntry] = []

    def fix(self, deck: Path, findings: list[Finding], *, out_dir: Path | None = None) -> Path:
        target_dir = out_dir or deck.parent
        target_dir.mkdir(parents=True, exist_ok=True)
        target = target_dir / FIXED_DECK_FILENAME
        shutil.copy2(deck, target)

        original_slide_count = len(Presentation(deck).slides)
        presentation = Presentation(target)
        log: list[FixLogEntry] = []
        touched: set[tuple[int, int]] = set()

        for finding in findings:
            entry = self._apply_one(finding, presentation, touched)
            if entry is not None:
                log.append(entry)

        presentation.save(target)

        if not self._passes_control_check(target, original_slide_count, touched):
            shutil.copy2(deck, target)
            for finding in findings:
                finding.auto_fixed = False
            log = [
                FixLogEntry(
                    finding_id=e.finding_id, rule=e.rule, status="skipped", reason="rolled_back"
                )
                for e in log
            ]

        self.last_log = log
        log_payload = [entry.model_dump(mode="json") for entry in log]
        (target_dir / "fix_log.json").write_text(
            json.dumps(log_payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return target

    def _apply_one(
        self,
        finding: Finding,
        presentation: Presentation,
        touched: set[tuple[int, int]],
    ) -> FixLogEntry | None:
        if not finding.auto_fixable:
            return None
        if finding.slide_num is None or finding.bbox is None:
            return FixLogEntry(
                finding_id=finding.id, status="skipped", reason="no slide_num/bbox"
            )
        if finding.slide_num < 1 or finding.slide_num > len(presentation.slides):
            return FixLogEntry(
                finding_id=finding.id, status="skipped", reason="slide_num out of range"
            )

        slide = presentation.slides[finding.slide_num - 1]
        shape = _find_matching_shape(
            slide, finding.bbox, presentation.slide_width, presentation.slide_height
        )
        if shape is None:
            return FixLogEntry(finding_id=finding.id, status="skipped", reason="no matching shape")

        rule = next((r for r in self._rules if r.applies_to(finding)), None)
        if rule is None:
            return FixLogEntry(finding_id=finding.id, status="skipped", reason="no applicable rule")

        result = rule.apply(shape, finding, slide)
        if result.applied:
            finding.auto_fixed = True
            touched.add((finding.slide_num, shape.shape_id))
            return FixLogEntry(finding_id=finding.id, rule=rule.name, status="applied")
        return FixLogEntry(
            finding_id=finding.id, rule=rule.name, status="skipped", reason=result.reason
        )

    @staticmethod
    def _passes_control_check(
        target: Path, original_slide_count: int, touched: set[tuple[int, int]]
    ) -> bool:
        presentation = Presentation(target)
        if len(presentation.slides) != original_slide_count:
            return False
        slide_width, slide_height = presentation.slide_width, presentation.slide_height
        for slide_num, shape_id in touched:
            slide = presentation.slides[slide_num - 1]
            shape = next((s for s in slide.shapes if s.shape_id == shape_id), None)
            if shape is None:
                return False
            if shape.left is None or shape.top is None:
                continue
            if shape.left < 0 or shape.top < 0:
                return False
            if shape.left + (shape.width or 0) > slide_width:
                return False
            if shape.top + (shape.height or 0) > slide_height:
                return False
        return True
