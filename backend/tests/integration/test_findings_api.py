"""httpx integration tests for findings flag / like / apply_fix."""

from __future__ import annotations

from collections.abc import AsyncIterator
from datetime import UTC, datetime, timedelta
from pathlib import Path
from uuid import UUID, uuid4

import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine
from sqlalchemy.pool import NullPool

import app.services.finding_service as finding_service_module
from app.config import get_settings
from app.db import get_db, reset_db_state
from app.deps import get_storage
from app.main import create_app
from app.models import Base, FileAsset, FileAssetKind, FindingRow, Review
from app.services.file_assets import save_file_asset
from app.services.storage import LocalStorage


@pytest.fixture
async def findings_client(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> AsyncIterator[tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage]]:
    db_path = tmp_path / "findings.db"
    storage_root = tmp_path / "storage"
    monkeypatch.setenv("DATABASE_URL", f"sqlite+aiosqlite:///{db_path}")
    monkeypatch.setenv("REDIS_URL", "redis://localhost:6379/0")
    monkeypatch.setenv("SECRET_KEY", "test-secret-key-for-jwt-32bytes!!")
    monkeypatch.setenv("LLM_API_KEY", "test-key")
    monkeypatch.setenv("SMTP_HOST", "")
    monkeypatch.setenv("STORAGE_ROOT", str(storage_root))
    monkeypatch.setenv("ENVIRONMENT", "test")
    get_settings.cache_clear()
    reset_db_state()

    engine = create_async_engine(f"sqlite+aiosqlite:///{db_path}", poolclass=NullPool)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    factory = async_sessionmaker(engine, expire_on_commit=False)
    storage = LocalStorage(storage_root)

    async def override_get_db() -> AsyncIterator[AsyncSession]:
        async with factory() as session:
            yield session

    app = create_app(get_settings())
    app.dependency_overrides[get_db] = override_get_db
    app.dependency_overrides[get_storage] = lambda: storage

    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        yield client, factory, storage

    app.dependency_overrides.clear()
    await engine.dispose()
    get_settings.cache_clear()
    reset_db_state()


async def _register(client: AsyncClient, email: str) -> str:
    reg = await client.post("/api/v1/auth/register", json={"email": email, "password": "password1"})
    assert reg.status_code == 201, reg.text
    access: str = reg.json()["access_token"]
    return access


async def _user_id(client: AsyncClient, access: str) -> UUID:
    me = await client.get("/api/v1/auth/me", headers={"Authorization": f"Bearer {access}"})
    return UUID(me.json()["id"])


def _tiny_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "tiny"
    run.font.size = Pt(10)
    presentation.save(path)


async def _seed_finding(
    factory: async_sessionmaker[AsyncSession],
    storage: LocalStorage,
    *,
    user_id: UUID,
    auto_fixable: bool = False,
    with_pptx: bool = False,
) -> UUID:
    async with factory() as session:
        review = Review(
            user_id=user_id,
            status="done",
            deck_filename="deck.pptx" if with_pptx else "deck.pdf",
            has_audio=False,
            has_data=False,
        )
        session.add(review)
        await session.flush()

        if with_pptx:
            pptx_path = Path(storage.root) / "seed.pptx"
            _tiny_pptx(pptx_path)
            data = pptx_path.read_bytes()
            await save_file_asset(
                storage,
                session,
                review_id=review.id,
                kind=FileAssetKind.DECK_ORIGINAL,
                filename="deck.pptx",
                data=data,
                expires_at=datetime.now(UTC) + timedelta(days=7),
            )
            await save_file_asset(
                storage,
                session,
                review_id=review.id,
                kind=FileAssetKind.FIXED_PPTX,
                filename="fixed.pptx",
                data=data,
                expires_at=datetime.now(UTC) + timedelta(days=7),
            )

        finding = FindingRow(
            review_id=review.id,
            slide_num=1,
            category="TYPOGRAPHY",
            severity="MINOR",
            title="Мелкий шрифт",
            description="Кегль меньше 14pt",
            fix_suggestion="Увеличьте шрифт",
            bbox={"x": 0.1, "y": 0.1, "w": 0.4, "h": 0.2},
            auto_fixable=auto_fixable,
            auto_fixed=False,
            source="slide_analyzer",
            user_flag=False,
            user_like=False,
        )
        session.add(finding)
        await session.commit()
        return finding.id


@pytest.mark.asyncio
async def test_owner_flags_finding(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, factory, storage = findings_client
    access = await _register(client, "flagger@example.com")
    user_id = await _user_id(client, access)
    finding_id = await _seed_finding(factory, storage, user_id=user_id)

    resp = await client.post(
        f"/api/v1/findings/{finding_id}/flag", headers={"Authorization": f"Bearer {access}"}
    )
    assert resp.status_code == 204

    async with factory() as check:
        row = await check.get(FindingRow, finding_id)
        assert row is not None
        assert row.user_flag is True
        assert row.user_like is False


@pytest.mark.asyncio
async def test_like_and_flag_are_mutually_exclusive(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, factory, storage = findings_client
    access = await _register(client, "voter@example.com")
    user_id = await _user_id(client, access)
    finding_id = await _seed_finding(factory, storage, user_id=user_id)
    headers = {"Authorization": f"Bearer {access}"}

    like = await client.post(f"/api/v1/findings/{finding_id}/like", headers=headers)
    assert like.status_code == 204
    async with factory() as check:
        row = await check.get(FindingRow, finding_id)
        assert row is not None
        assert row.user_like is True
        assert row.user_flag is False

    flag = await client.post(f"/api/v1/findings/{finding_id}/flag", headers=headers)
    assert flag.status_code == 204
    async with factory() as check:
        row = await check.get(FindingRow, finding_id)
        assert row is not None
        assert row.user_flag is True
        assert row.user_like is False


@pytest.mark.asyncio
async def test_flagging_is_idempotent(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, factory, storage = findings_client
    access = await _register(client, "twice@example.com")
    user_id = await _user_id(client, access)
    finding_id = await _seed_finding(factory, storage, user_id=user_id)
    headers = {"Authorization": f"Bearer {access}"}

    first = await client.post(f"/api/v1/findings/{finding_id}/flag", headers=headers)
    second = await client.post(f"/api/v1/findings/{finding_id}/flag", headers=headers)
    assert first.status_code == 204
    assert second.status_code == 204


@pytest.mark.asyncio
async def test_non_owner_gets_404_on_like_and_apply(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, factory, storage = findings_client
    owner_access = await _register(client, "flagowner@example.com")
    owner_id = await _user_id(client, owner_access)
    other_access = await _register(client, "flagother@example.com")
    finding_id = await _seed_finding(
        factory, storage, user_id=owner_id, auto_fixable=True, with_pptx=True
    )

    like = await client.post(
        f"/api/v1/findings/{finding_id}/like",
        headers={"Authorization": f"Bearer {other_access}"},
    )
    apply = await client.post(
        f"/api/v1/findings/{finding_id}/apply_fix",
        headers={"Authorization": f"Bearer {other_access}"},
    )
    assert like.status_code == 404
    assert apply.status_code == 404


@pytest.mark.asyncio
async def test_unknown_finding_gets_404(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, _factory, _storage = findings_client
    access = await _register(client, "flagunknown@example.com")

    resp = await client.post(
        f"/api/v1/findings/{uuid4()}/flag", headers={"Authorization": f"Bearer {access}"}
    )
    assert resp.status_code == 404


@pytest.mark.asyncio
async def test_apply_fix_requires_auto_fixable(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, factory, storage = findings_client
    access = await _register(client, "noapply@example.com")
    user_id = await _user_id(client, access)
    finding_id = await _seed_finding(
        factory, storage, user_id=user_id, auto_fixable=False, with_pptx=True
    )

    resp = await client.post(
        f"/api/v1/findings/{finding_id}/apply_fix",
        headers={"Authorization": f"Bearer {access}"},
    )
    assert resp.status_code == 404


@pytest.mark.asyncio
async def test_apply_fix_sets_auto_fixed_and_is_idempotent(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
) -> None:
    client, factory, storage = findings_client
    access = await _register(client, "applier@example.com")
    user_id = await _user_id(client, access)
    finding_id = await _seed_finding(
        factory, storage, user_id=user_id, auto_fixable=True, with_pptx=True
    )
    headers = {"Authorization": f"Bearer {access}"}

    first = await client.post(f"/api/v1/findings/{finding_id}/apply_fix", headers=headers)
    assert first.status_code == 204
    async with factory() as check:
        row = await check.get(FindingRow, finding_id)
        assert row is not None
        assert row.auto_fixed is True

        fixed_asset = (
            await check.execute(
                select(FileAsset).where(
                    FileAsset.review_id == row.review_id,
                    FileAsset.kind == FileAssetKind.FIXED_PPTX.value,
                )
            )
        ).scalar_one()
        assert fixed_asset.version == 2
        assert fixed_asset.storage_path.endswith("_deck_исправленная_версия№2.pptx")

    second = await client.post(f"/api/v1/findings/{finding_id}/apply_fix", headers=headers)
    assert second.status_code == 204
    async with factory() as check:
        fixed_asset = (
            await check.execute(
                select(FileAsset).where(
                    FileAsset.review_id == row.review_id,
                    FileAsset.kind == FileAssetKind.FIXED_PPTX.value,
                )
            )
        ).scalar_one()
        assert fixed_asset.version == 2, "second (idempotent) apply must not bump the version again"


@pytest.mark.asyncio
async def test_langfuse_score_called_with_expected_fields(
    findings_client: tuple[AsyncClient, async_sessionmaker[AsyncSession], LocalStorage],
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    client, factory, storage = findings_client
    access = await _register(client, "langfusecalled@example.com")
    user_id = await _user_id(client, access)
    finding_id = await _seed_finding(factory, storage, user_id=user_id)

    flag_calls: list[dict[str, object]] = []
    like_calls: list[dict[str, object]] = []

    def fake_flag(settings: object, **kwargs: object) -> None:
        flag_calls.append(kwargs)

    def fake_like(settings: object, **kwargs: object) -> None:
        like_calls.append(kwargs)

    monkeypatch.setattr(finding_service_module, "score_finding_flag", fake_flag)
    monkeypatch.setattr(finding_service_module, "score_finding_like", fake_like)

    headers = {"Authorization": f"Bearer {access}"}
    assert (
        await client.post(f"/api/v1/findings/{finding_id}/flag", headers=headers)
    ).status_code == 204
    assert (
        await client.post(f"/api/v1/findings/{finding_id}/like", headers=headers)
    ).status_code == 204

    assert len(flag_calls) == 1
    assert flag_calls[0]["finding_id"] == finding_id
    assert len(like_calls) == 1
    assert like_calls[0]["category"] == "TYPOGRAPHY"
