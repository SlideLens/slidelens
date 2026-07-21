"""Unit tests for ``core.ingest`` (DeckIngestor, AudioExtractor)."""

from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from core.constants import WAV_16K_MONO_BYTES_PER_SECOND
from core.context import ReviewContext
from core.ingest import AudioExtractor, DeckIngestor
from core.ingest_errors import (
    AudioTooLongError,
    CorruptedFileError,
    DeckTooLargeError,
    EmptyDeckError,
    NoAudioTrackError,
    RenderTimeoutError,
    UnsupportedDeckFormatError,
)


def _make_pptx(path: Path, slide_texts: list[str]) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = text
    prs.save(path)


def _make_pdf(path: Path, n_pages: int) -> None:
    images = [Image.new("RGB", (200, 150), "white") for _ in range(n_pages)]
    images[0].save(path, save_all=True, append_images=images[1:])


def _make_ctx(workdir: Path) -> ReviewContext:
    return ReviewContext(workdir=workdir)


@pytest.mark.asyncio
async def test_ingest_pdf_renders_numbered_pngs(tmp_path: Path) -> None:
    pdf_path = tmp_path / "deck.pdf"
    _make_pdf(pdf_path, 3)
    workdir = tmp_path / "work"

    pngs = await DeckIngestor().ingest(pdf_path, workdir)

    assert [p.name for p in pngs] == ["slide_001.png", "slide_002.png", "slide_003.png"]
    assert all(p.is_file() for p in pngs)


@pytest.mark.asyncio
async def test_ingest_pptx_extracts_text_and_registers_ctx(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, ["Hello slide one", "Second slide"])
    workdir = tmp_path / "work"
    ctx = _make_ctx(workdir)

    async def fake_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        outdir = Path(cmd[cmd.index("--outdir") + 1])
        pdf_path = outdir / f"{pptx_path.stem}.pdf"
        _make_pdf(pdf_path, 2)
        return 0, b"", b""

    ingestor = DeckIngestor(run_subprocess=fake_run_subprocess)
    pngs = await ingestor.ingest(pptx_path, workdir, ctx)

    assert len(pngs) == 2
    assert ctx.meta["slide_texts"][1] == "Hello slide one"
    assert ctx.meta["slide_texts"][2] == "Second slide"
    assert ctx.slide_pngs[1] == pngs[0]
    assert ctx.slide_pngs[2] == pngs[1]


@pytest.mark.asyncio
async def test_ingest_oversized_deck_rejected_before_render(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr("core.ingest.MAX_DECK_SIZE_MB", 0)
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, ["only slide"])

    called = False

    async def fail_if_called(*args: object, **kwargs: object) -> tuple[int, bytes, bytes]:
        nonlocal called
        called = True
        return 0, b"", b""

    ingestor = DeckIngestor(run_subprocess=fail_if_called)
    with pytest.raises(DeckTooLargeError):
        await ingestor.ingest(pptx_path, tmp_path / "work")
    assert called is False


@pytest.mark.asyncio
async def test_ingest_empty_pptx_raises(tmp_path: Path) -> None:
    pptx_path = tmp_path / "empty.pptx"
    _make_pptx(pptx_path, [])

    with pytest.raises(EmptyDeckError):
        await DeckIngestor().ingest(pptx_path, tmp_path / "work")


@pytest.mark.asyncio
async def test_ingest_corrupted_pptx_raises(tmp_path: Path) -> None:
    bad_path = tmp_path / "bad.pptx"
    bad_path.write_bytes(b"not a real pptx file at all")

    with pytest.raises(CorruptedFileError):
        await DeckIngestor().ingest(bad_path, tmp_path / "work")


@pytest.mark.asyncio
async def test_ingest_unsupported_format_raises(tmp_path: Path) -> None:
    ppt_path = tmp_path / "legacy.ppt"
    ppt_path.write_bytes(b"whatever")

    with pytest.raises(UnsupportedDeckFormatError):
        await DeckIngestor().ingest(ppt_path, tmp_path / "work")


@pytest.mark.asyncio
async def test_ingest_garbage_pdf_raises_corrupted(tmp_path: Path) -> None:
    fake_pdf = tmp_path / "not_a_pdf.pdf"
    fake_pdf.write_text("this is definitely not a pdf")

    with pytest.raises(CorruptedFileError):
        await DeckIngestor().ingest(fake_pdf, tmp_path / "work")


@pytest.mark.asyncio
async def test_soffice_retry_then_succeed(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, ["slide one"])
    workdir = tmp_path / "work"
    calls = {"n": 0}

    async def flaky_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        calls["n"] += 1
        if calls["n"] == 1:
            return 1, b"", b"transient failure"
        outdir = Path(cmd[cmd.index("--outdir") + 1])
        _make_pdf(outdir / f"{pptx_path.stem}.pdf", 1)
        return 0, b"", b""

    pngs = await DeckIngestor(run_subprocess=flaky_run_subprocess).ingest(pptx_path, workdir)

    assert calls["n"] == 2
    assert len(pngs) == 1


@pytest.mark.asyncio
async def test_soffice_timeout_raises_render_timeout(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path, ["slide one"])

    async def always_times_out(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        raise TimeoutError

    with pytest.raises(RenderTimeoutError):
        await DeckIngestor(run_subprocess=always_times_out).ingest(pptx_path, tmp_path / "work")


@pytest.mark.asyncio
async def test_audio_extract_success(tmp_path: Path) -> None:
    media_path = tmp_path / "clip.mp4"
    media_path.write_bytes(b"fake media bytes")

    async def fake_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        out_path = Path(cmd[-1])
        out_path.write_bytes(b"RIFF....WAVEfmt ")
        return 0, b"", b""

    wav_path = await AudioExtractor(run_subprocess=fake_run_subprocess).extract(
        media_path, tmp_path / "work"
    )

    assert wav_path.is_file()
    assert wav_path.suffix == ".wav"


def _wav_of_minutes(minutes: float) -> bytes:
    """WAV нужной длительности: 16 кГц моно 16 бит = 32000 байт/с."""
    return b"R" * int(minutes * 60 * WAV_16K_MONO_BYTES_PER_SECOND)


@pytest.mark.asyncio
async def test_audio_longer_than_the_cap_is_rejected(tmp_path: Path) -> None:
    """Транскрипция тарифицируется поминутно — без потолка Разбор стоит сколько угодно."""
    media_path = tmp_path / "long.mp4"
    media_path.write_bytes(b"fake media bytes")

    async def fake_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        Path(cmd[-1]).write_bytes(_wav_of_minutes(31))
        return 0, b"", b""

    with pytest.raises(AudioTooLongError):
        await AudioExtractor(run_subprocess=fake_run_subprocess, max_minutes=30).extract(
            media_path, tmp_path / "work"
        )

    # Огромный WAV не должен остаться лежать на диске после отказа.
    assert not (tmp_path / "work" / "long.wav").exists()


@pytest.mark.asyncio
async def test_audio_just_under_the_cap_is_accepted(tmp_path: Path) -> None:
    media_path = tmp_path / "ok.mp4"
    media_path.write_bytes(b"fake media bytes")

    async def fake_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        Path(cmd[-1]).write_bytes(_wav_of_minutes(29))
        return 0, b"", b""

    wav_path = await AudioExtractor(run_subprocess=fake_run_subprocess, max_minutes=30).extract(
        media_path, tmp_path / "work"
    )

    assert wav_path.is_file()


@pytest.mark.asyncio
async def test_audio_extract_no_audio_track(tmp_path: Path) -> None:
    media_path = tmp_path / "silent.mp4"
    media_path.write_bytes(b"fake media bytes")

    async def fake_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        return 1, b"", b"Output file does not contain any stream"

    with pytest.raises(NoAudioTrackError):
        await AudioExtractor(run_subprocess=fake_run_subprocess).extract(
            media_path, tmp_path / "work"
        )


@pytest.mark.asyncio
async def test_audio_extract_generic_failure_is_corrupted(tmp_path: Path) -> None:
    media_path = tmp_path / "bad.mp4"
    media_path.write_bytes(b"fake media bytes")

    async def fake_run_subprocess(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        return 1, b"", b"Invalid data found when processing input"

    with pytest.raises(CorruptedFileError):
        await AudioExtractor(run_subprocess=fake_run_subprocess).extract(
            media_path, tmp_path / "work"
        )


@pytest.mark.asyncio
async def test_audio_extract_timeout_raises_render_timeout(tmp_path: Path) -> None:
    media_path = tmp_path / "clip.mp4"
    media_path.write_bytes(b"fake media bytes")

    async def always_times_out(cmd: list[str], *, timeout: float) -> tuple[int, bytes, bytes]:
        raise TimeoutError

    with pytest.raises(RenderTimeoutError):
        await AudioExtractor(run_subprocess=always_times_out).extract(
            media_path, tmp_path / "work"
        )
