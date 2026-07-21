"""Unit tests for ``core.fix`` (FixRules + PptxFixer)."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from core.fix import FixRule, PptxFixer
from core.fix_schemas import FixResult
from core.schemas import BBox, Category, Finding, Severity


def _slide_box_bbox(prs: Presentation, box) -> BBox:
    return BBox(
        x=box.left / prs.slide_width,
        y=box.top / prs.slide_height,
        w=box.width / prs.slide_width,
        h=box.height / prs.slide_height,
    )


def _finding(**overrides: object) -> Finding:
    payload: dict[str, object] = {
        "slide_num": 1,
        "category": Category.TYPOGRAPHY,
        "severity": Severity.MINOR,
        "title": "t",
        "description": "d",
        "fix_suggestion": "f",
        "auto_fixable": True,
    }
    payload.update(overrides)
    return Finding(**payload)


def test_min_font_size_rule_raises_small_font(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "small text"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    bbox = _slide_box_bbox(prs, box)
    prs.save(deck_path)

    finding = _finding(
        title="Мелкий шрифт",
        description="Шрифт заголовка слишком мелкий",
        bbox=bbox,
    )

    out_path = PptxFixer().fix(deck_path, [finding], out_dir=tmp_path)

    assert out_path.name == "fixed.pptx"
    fixed = Presentation(out_path)
    run = fixed.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(14)
    assert finding.auto_fixed is True


def test_contrast_rule_recolors_low_contrast_text(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.text_frame.text = "low contrast"
    box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(230, 230, 230)
    bbox = _slide_box_bbox(prs, box)
    prs.save(deck_path)

    finding = _finding(
        category=Category.READABILITY,
        title="Низкий контраст текста",
        description="Текст низкого контраста на светлом фоне",
        bbox=bbox,
    )

    out_path = PptxFixer().fix(deck_path, [finding], out_dir=tmp_path)

    fixed = Presentation(out_path)
    run = fixed.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert tuple(run.font.color.rgb) == (0, 0, 0)
    assert finding.auto_fixed is True


def test_alignment_rule_snaps_near_aligned_edge(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    anchor = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    anchor.text_frame.text = "anchor"
    moved = slide.shapes.add_textbox(Inches(1) + 20000, Inches(3), Inches(2), Inches(1))
    moved.text_frame.text = "moved, almost aligned"
    bbox = _slide_box_bbox(prs, moved)
    prs.save(deck_path)

    finding = _finding(
        category=Category.HIERARCHY,
        title="Блок почти выровнен",
        description="Блок почти выровнен по левому краю с соседним",
        bbox=bbox,
    )

    out_path = PptxFixer().fix(deck_path, [finding], out_dir=tmp_path)

    fixed = Presentation(out_path)
    lefts = {shape.left for shape in fixed.slides[0].shapes}
    assert len(lefts) == 1
    assert finding.auto_fixed is True


def test_non_auto_fixable_finding_is_untouched_and_unlogged(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "small text"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    bbox = _slide_box_bbox(prs, box)
    prs.save(deck_path)

    finding = _finding(bbox=bbox, auto_fixable=False)
    fixer = PptxFixer()

    out_path = fixer.fix(deck_path, [finding], out_dir=tmp_path)

    fixed = Presentation(out_path)
    run = fixed.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(8)
    assert finding.auto_fixed is False
    assert fixer.last_log == []


def test_no_matching_shape_is_logged_and_does_not_raise(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(deck_path)

    finding = _finding(bbox=BBox(x=0.1, y=0.1, w=0.1, h=0.1))
    fixer = PptxFixer()

    out_path = fixer.fix(deck_path, [finding], out_dir=tmp_path)

    assert out_path.is_file()
    assert fixer.last_log[0].status == "skipped"
    assert fixer.last_log[0].reason == "no matching shape"


def test_no_applicable_rule_is_logged(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "hello"
    bbox = _slide_box_bbox(prs, box)
    prs.save(deck_path)

    finding = _finding(
        category=Category.NARRATIVE, title="Провал нарратива", description="нет CTA", bbox=bbox
    )
    fixer = PptxFixer()

    fixer.fix(deck_path, [finding], out_dir=tmp_path)

    assert fixer.last_log[0].status == "skipped"
    assert fixer.last_log[0].reason == "no applicable rule"


class _BreakingRule(FixRule):
    name = "breaking"

    def applies_to(self, finding: Finding) -> bool:
        return True

    def apply(self, shape: object, finding: Finding, slide: object) -> FixResult:
        shape.left = -shape.width  # push it off-slide
        return FixResult(applied=True)


def test_control_check_rolls_back_out_of_bounds_edit(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "hello"
    bbox = _slide_box_bbox(prs, box)
    prs.save(deck_path)

    finding = _finding(bbox=bbox)
    fixer = PptxFixer(rules=[_BreakingRule()])

    out_path = fixer.fix(deck_path, [finding], out_dir=tmp_path)

    fixed = Presentation(out_path)
    assert fixed.slides[0].shapes[0].left == box.left  # restored to original
    assert finding.auto_fixed is False
    assert fixer.last_log[0].status == "skipped"
    assert fixer.last_log[0].reason == "rolled_back"


def test_output_always_opens_with_no_findings(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(deck_path)

    out_path = PptxFixer().fix(deck_path, [], out_dir=tmp_path)

    assert len(Presentation(out_path).slides) == 1


@pytest.mark.parametrize(
    ("title", "description", "expected"),
    [
        ("Мелкий шрифт", "маленький кегль", True),
        ("Не то", "совсем другая проблема", False),
    ],
)
def test_min_font_size_rule_applies_to_keyword_match(
    title: str, description: str, expected: bool
) -> None:
    from core.fix import MinFontSizeRule

    finding = _finding(title=title, description=description)
    assert MinFontSizeRule().applies_to(finding) is expected
