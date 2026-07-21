"""One-off generator for the golden deck fixtures under ``decks/``.

Run with ``uv run python tests/golden/generate_decks.py`` to (re)create them.
Kept for reproducibility — the golden set is checked in as binary ``.pptx``
files, this script is how they were made and how to regenerate them if the
expected findings change.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt

DECKS_DIR = Path(__file__).resolve().parent / "decks"


def _add_text_slide(prs: Presentation, title: str, body: str, *, font_pt: int = 24) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    body_box.text_frame.text = body
    for run in body_box.text_frame.paragraphs[0].runs:
        run.font.size = Pt(font_pt)


def _bar_chart_image(path: Path, *, values: list[int], truncated_axis: bool) -> None:
    img = Image.new("RGB", (800, 500), "white")
    draw = ImageDraw.Draw(img)
    max_v = max(values)
    baseline = 450
    floor_v = max_v * 0.7 if truncated_axis else 0
    for i, v in enumerate(values):
        frac = (v - floor_v) / (max_v - floor_v) if max_v != floor_v else 1.0
        height = int(frac * 350)
        x0 = 100 + i * 150
        draw.rectangle([x0, baseline - height, x0 + 100, baseline], fill="steelblue")
        draw.text((x0, baseline + 10), str(v), fill="black")
    draw.line([(60, baseline), (760, baseline)], fill="black", width=2)
    img.save(path, "PNG")


def _pie_chart_image(path: Path, *, shares: list[int]) -> None:
    img = Image.new("RGB", (500, 500), "white")
    draw = ImageDraw.Draw(img)
    colors = ["steelblue", "indianred", "seagreen", "goldenrod"]
    start = 0.0
    total = sum(shares)
    for i, share in enumerate(shares):
        extent = share / total * 360 if total else 0
        draw.pieslice([50, 50, 450, 450], start, start + extent, fill=colors[i % len(colors)])
        start += extent
    img.save(path, "PNG")


def build_plain_ru_1() -> None:
    prs = Presentation()
    _add_text_slide(prs, "Проблема", "Малый бизнес теряет 20% выручки на ручном учёте.")
    _add_text_slide(prs, "Решение", "SlideLens автоматизирует разбор презентаций перед питчем.")
    _add_text_slide(prs, "Доказательства", "50 пилотных пользователей, NPS 62.")
    _add_text_slide(prs, "Призыв к действию", "Инвестируйте 200 000$ на раунде Seed.")
    prs.save(DECKS_DIR / "plain_ru_1.pptx")


def build_plain_ru_2() -> None:
    prs = Presentation()
    _add_text_slide(prs, "О команде", "Три сооснователя с опытом в EdTech и SaaS.")
    _add_text_slide(prs, "Рынок", "Объём рынка — 4 млрд рублей в РФ.", font_pt=10)
    _add_text_slide(prs, "Дорожная карта", "Q1 — MVP, Q2 — пилоты, Q3 — масштабирование.")
    prs.save(DECKS_DIR / "plain_ru_2.pptx")


def build_chart_bar() -> None:
    prs = Presentation()
    _add_text_slide(prs, "Рост выручки", "Выручка стабильно растёт квартал к кварталу.")
    chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_path = DECKS_DIR / "_chart_bar_truncated.png"
    _bar_chart_image(chart_path, values=[92, 96, 100], truncated_axis=True)
    chart_slide.shapes.add_picture(str(chart_path), Inches(1), Inches(1), Inches(6), Inches(4))
    _add_text_slide(prs, "Итог", "Продолжаем тот же темп роста в следующем году.")
    prs.save(DECKS_DIR / "chart_bar.pptx")
    chart_path.unlink()


def build_chart_pie() -> None:
    prs = Presentation()
    _add_text_slide(prs, "Структура выручки", "Три сегмента формируют портфель продукта.")
    chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_path = DECKS_DIR / "_chart_pie.png"
    _pie_chart_image(chart_path, shares=[40, 35, 17])
    chart_slide.shapes.add_picture(str(chart_path), Inches(2), Inches(1), Inches(4), Inches(4))
    _add_text_slide(prs, "Итог", "Диверсифицированный портфель снижает риски.")
    prs.save(DECKS_DIR / "chart_pie.pptx")
    chart_path.unlink()


def build_bad_deck() -> None:
    prs = Presentation()
    _add_text_slide(prs, "Введение", "Компания основана в 2023 году в Москве.")
    _add_text_slide(
        prs,
        "Мелкий текст",
        "Этот текст набран слишком мелким шрифтом для проекции в зале.",
        font_pt=8,
    )
    dup_slide_1 = prs.slides.add_slide(prs.slide_layouts[6])
    dup_box = dup_slide_1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    dup_box.text_frame.text = "Наша команда — эксперты рынка с 10-летним опытом."
    dup_slide_2 = prs.slides.add_slide(prs.slide_layouts[6])
    dup_box2 = dup_slide_2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    dup_box2.text_frame.text = "Наша команда — эксперты рынка с 10-летним опытом."
    prs.save(DECKS_DIR / "bad_deck.pptx")


def main() -> None:
    DECKS_DIR.mkdir(parents=True, exist_ok=True)
    build_plain_ru_1()
    build_plain_ru_2()
    build_chart_bar()
    build_chart_pie()
    build_bad_deck()


if __name__ == "__main__":
    main()
